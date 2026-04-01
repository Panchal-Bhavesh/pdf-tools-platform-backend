import sys
import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

try:
    if len(sys.argv) != 3:
        print("ERROR: Incorrect number of arguments")
        sys.exit(1)

    pdf_path = sys.argv[1]
    pptx_path = sys.argv[2]

    print(f"Converting {pdf_path} to {pptx_path}")

    if not os.path.exists(pdf_path):
        print(f"ERROR: PDF file does not exist: {pdf_path}")
        sys.exit(1)

    # Convert PDF pages to images
    images = convert_from_path(pdf_path)

    if not images:
        print("ERROR: No pages found in PDF")
        sys.exit(1)

    # Create a PowerPoint presentation
    prs = Presentation()

    for i, image in enumerate(images):
        # Add a slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

        # Save image temporarily
        img_path = f"temp_slide_{i}.png"
        image.save(img_path, 'PNG')

        # Add image to slide
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

        # Clean up temp image
        os.remove(img_path)

    # Save the presentation
    prs.save(pptx_path)

    if os.path.exists(pptx_path):
        print("SUCCESS")
    else:
        print("ERROR: Output file was not created")
        sys.exit(1)

except Exception as e:
    print(f"ERROR: {str(e)}")
    sys.exit(1)